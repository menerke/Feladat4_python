import logging

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm

#from presentation_io import PresentationIO

class PresentationGenerator:
    """The generator class.

    It creates different type of slides. The available types: Title, Text, Image, List, Plot.
    """
    def __init__(self, outputFileName,templateFileName):
        """Initialization.
        
        Set up the logger. Give the template file of the presentation.

        Parameters
        ----------
        outputFileName : str
            The name of the output file.
        templateFileName : str
            The name of the template file. 
        Raises
        ----------
        TypeError
            If the outputFileName or templateFileName is not a string.

        ValueError
            If the outputFileName has not .pptx extension or templateFileName has not .template extension.
        """
        self._logger = logging.getLogger(self.__class__.__name__)
        if not type(outputFileName) or not type(templateFileName) is str:
            self._logger.error('Name of the file must be string.')
            raise TypeError
        else:
            if outputFileName.endswith('.pptx') or templateFileName.endswith('.template'):
                self._outputFileName = outputFileName
            else:
                self._logger.error('The file extension is wrong.')
                raise ValueError

        self._presentation = Presentation(templateFileName)
        self.outputFileName = outputFileName

    def addTitle(self, layout, title, subTitle):
        """Generate the Title slide. 

        The slide contains the title and the subtitle. The function should select the layout first, then add a new slide and write the title and subtitle text.
        
        Parameters
        ----------
        layout : int
            The number of the layout to be selected. 
        title : str
            The string contains the title text. 
        subTitle : str
            The string contains the subtitle text. 

        Returns
        -------
        bool
            True if the slide generation is successful.

        Raises
        ----------
        SystemError
            If the slide generation is not successful. 
        """
        slideLayout = self._presentation.slide_layouts[layout]
        try:
            slide = self._presentation.slides.add_slide(slideLayout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = subTitle
            self._logger.info("Title page is added ({0}, {1})".format(title, subTitle))
            return True
        except:
            self._logger.error("The slide generation is not successful.")
            raise SystemError

    def addText(self, layout, title, text):
        """Generate the Text slide. 
        
        The slide contains a title and a longer text. It should select the layout first, then add a new slide and write the title and the long text.
        
        Parameters
        ----------
        layout : int
            The number of the layout to be selected. 
        title : str
            The string contains the title text. 
        text : str
            The string contains the long text. 

        Returns
        -------
        bool
            True if the slide generation is successful.

        Raises
        ----------
        SystemError
            If the slide generation is not successful. 
        
        """
        slideLayout = self._presentation.slide_layouts[layout]
        try:
            slide = self._presentation.slides.add_slide(slideLayout)
            titleShape = slide.shapes.title
            titleShape.text = title
            left = Cm(3.5)
            top = Cm(3.0)
            height = Cm(14.5)
            width = Cm(30.0)
            textBox = slide.shapes.add_textbox(left, top, height, width)
            textFrame = textBox.text_frame
            textFrame.text = text
            self._logger.info("Title page is added ({0}, {1})".format(title, text))
            return True
        except:
            self._logger.error("The slide generation is not successful.")
            raise SystemError

    def addList(self, layout, title, levels, text):
        """Generate the List slide. 
        
        The slide contains a title and a list. The number of the level is also an input. 
        It should select the layout first, then add a new slide and write the title and the text on the appropriate level.
        
        Parameters
        ----------
        layout : int
            The number of the layout to be selected. 
        title : str
            The string contains the title text. 
        levels: numpy array
            The numpy array contains the numbers to select the level of the list. 
        text : numpy array
            The numpy array contains the the text line by line (as a string). 

        Returns
        -------
        bool
            True if the slide generation is successful.

        Raises
        ----------
        SystemError
            If the slide generation is not successful. 
        """

        slideLayout = self._presentation.slide_layouts[layout]
        try:
            slide = self._presentation.slides.add_slide(slideLayout)
            shapes = slide.shapes
            titleShape = shapes.title
            bodyShape = shapes.placeholders[1]
            titleShape.text = title
            textFrame = bodyShape.text_frame
            lines = [levels,text]
            for i in range(len(levels)):
                paragraph = textFrame.add_paragraph()
                paragraph.level = lines[0][i]
                paragraph.text = lines[1][i]
            self._logger.info("List page is added.")
        except:
            self._logger.error("The slide generation is not successful.")
            raise SystemError
    
    def addImage(self, layout, title, fileName):
        """Generate the Image slide. 
            
            The slide contains a title and an image. It should select the layout first, then add a new slide and write the title and add the image.
            
            Parameters
            ----------
            layout : int
                The number of the layout to be selected. 
            title : str
                The string contains the title text. 
            fileName: str
                The name of the image file.

            Returns
            -------
            bool
                True if the slide generation is successful.

            Raises
            ----------
            SystemError
                If the slide generation is not successful. 
        """
        slideLayout = self._presentation.slide_layouts[layout]
        try:
            slide = self._presentation.slides.add_slide(slideLayout)
            titleShape = slide.shapes.title
            titleShape.text = title
            left = Cm(3.5)
            top = Cm(3.0)
            slide.shapes.add_picture(fileName, left, top)
            self._logger.info("Image page is added ({0}, {1})".format(title, fileName))
        except:
            self._logger.error("The slide generation is not successful.")
            raise SystemError

    def addPlot(self, layout, title, plotName):
        """Generate the Plot slide.
        
        The slide contains a title and an image named plotName. It should select the layout first, then add a new slide and write the title and add the image.
                    
        Parameters
        ----------
        layout : int
            The number of the layout to be selected. 
        title : str
            The string contains the title text. 
        plotName: str
            The name of the image file.

        Returns
        -------
        bool
            True if the slide generation is successful, False otherwise.

        Raises
        ----------
        SystemError
            If the slide generation is not successful. 

        """
        slideLayout = self._presentation.slide_layouts[layout]
        try:
            slide = self._presentation.slides.add_slide(slideLayout)
            titleShape = slide.shapes.title
            titleShape.text = title
            left = Cm(3.5)
            top = Cm(3.0)
            slide.shapes.add_picture(plotName, left, top)
            self._logger.info("Plot page is added ({0}, {1})".format(title, plotName))
        except:
            raise SystemError

    def finalize(self):
        """Save the created pptx.
    
        Returns
        -------
        bool
            True if the slide generation is successful.

        Raises
        ----------
        IOError
            If the finalization is not successfull, so the pptx file can not be written. 
        """
        try:
            self._presentation.save(self._outputFileName)
            return True
        except:
            raise IOError